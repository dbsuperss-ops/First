"""원본 슬라이드를 대상 프레젠테이션으로 복사합니다. 이미지/미디어 관계도 함께 복사합니다."""

import copy
from pptx import Presentation
from pptx.opc.part import Part
from pptx.opc.packuri import PackURI

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_R_EMBED = f"{{{_R_NS}}}embed"
_R_LINK = f"{{{_R_NS}}}link"
_R_ID = f"{{{_R_NS}}}id"

_CT_TO_EXT = {
    "image/jpeg": "jpg",
    "image/png": "png",
    "image/gif": "gif",
    "image/bmp": "bmp",
    "image/x-wmf": "wmf",
    "image/x-emf": "emf",
    "image/tiff": "tif",
}


def _get_ext(content_type: str) -> str:
    return _CT_TO_EXT.get(content_type, content_type.split("/")[-1])


def _copy_image_rels(src_slide, dest_slide) -> dict:
    """이미지 관계를 복사하고 rId 매핑을 반환합니다."""
    rId_map = {}
    counter = 0

    for src_rId, rel in src_slide.part.rels.items():
        if rel.is_external or "image" not in rel.reltype:
            continue

        src_img = rel._target
        ext = _get_ext(src_img.content_type)
        counter += 1

        # 고유한 파트 이름 생성 (슬라이드 객체 id로 충돌 방지)
        part_name = PackURI(f"/ppt/media/cp_{id(dest_slide)}_{counter}.{ext}")
        new_part = Part(part_name, src_img.content_type, src_img.blob)
        new_rId = dest_slide.part.relate_to(new_part, rel.reltype)
        rId_map[src_rId] = new_rId

    return rId_map


def _update_rids(tree, rId_map: dict):
    """XML 트리 내 r:embed / r:link / r:id 속성값을 새 rId로 교체합니다."""
    if not rId_map:
        return
    for el in tree.iter():
        for attr in (_R_EMBED, _R_LINK, _R_ID):
            val = el.get(attr)
            if val and val in rId_map:
                el.set(attr, rId_map[val])


def _copy_background(src_slide, dest_slide):
    """슬라이드 배경을 복사합니다."""
    src_bg = src_slide.background
    dest_bg = dest_slide.background

    if src_bg.fill.type is None:
        return

    src_elem = src_bg._element
    dest_elem = dest_bg._element
    parent = dest_elem.getparent()
    if parent is not None:
        parent.replace(dest_elem, copy.deepcopy(src_elem))


def copy_slide(src_slide, src_prs: Presentation, dest_prs: Presentation):
    """슬라이드를 원본에서 대상 프레젠테이션으로 복사합니다."""
    # 빈 슬라이드 추가 (blank layout = index 6)
    blank_layout = dest_prs.slide_layouts[6]
    dest_slide = dest_prs.slides.add_slide(blank_layout)

    # 자동 생성된 placeholder 요소 제거
    sp_tree = dest_slide.shapes._spTree
    for el in list(sp_tree):
        sp_tree.remove(el)

    # 이미지 관계 복사 및 rId 매핑 생성
    rId_map = _copy_image_rels(src_slide, dest_slide)

    # 원본 도형 트리를 deep copy 후 rId 갱신
    for src_el in src_slide.shapes._spTree:
        cloned = copy.deepcopy(src_el)
        _update_rids(cloned, rId_map)
        sp_tree.append(cloned)

    # 배경 복사
    _copy_background(src_slide, dest_slide)

    return dest_slide
