"""도형에 통일 서식을 적용합니다."""

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from config_loader import FormatConfig

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _apply_to_run(run, fmt: FormatConfig):
    """텍스트 런(run) 단위 서식 적용."""
    run.font.name = fmt.font_name
    run.font.size = Pt(fmt.font_size_pt)
    run.font.bold = fmt.bold
    if fmt.color_hex:
        run.font.color.rgb = _hex_to_rgb(fmt.color_hex)
    if fmt.char_spacing != 0.0:
        # OOXML spc 단위: 1/100 pt
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(fmt.char_spacing * 100)))


def _apply_to_paragraph(para, fmt: FormatConfig):
    """단락(paragraph) 단위 서식 적용."""
    if fmt.align in _ALIGN_MAP:
        para.alignment = _ALIGN_MAP[fmt.align]
    if fmt.line_spacing:
        para.line_spacing = fmt.line_spacing
    for run in para.runs:
        _apply_to_run(run, fmt)


def _apply_to_text_frame(tf, fmt: FormatConfig):
    for para in tf.paragraphs:
        _apply_to_paragraph(para, fmt)


def _apply_to_table(table, fmt: FormatConfig):
    for row in table.rows:
        for cell in row.cells:
            _apply_to_text_frame(cell.text_frame, fmt)


def apply_format(shape, fmt: FormatConfig):
    """도형 유형에 따라 서식을 적용합니다. 텍스트 없는 도형(이미지 등)은 건너뜁니다."""
    try:
        # 텍스트 프레임
        if shape.has_text_frame:
            _apply_to_text_frame(shape.text_frame, fmt)

        # 표(Table)
        if shape.has_table:
            _apply_to_table(shape.table, fmt)

        # 그룹: 자식 도형에 재귀 적용
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                apply_format(child, fmt)

        # 배경색 (설정된 경우에만)
        if fmt.bg_color_hex and shape.has_text_frame:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(fmt.bg_color_hex)

        # 테두리 (설정된 경우에만)
        if fmt.border_color_hex and fmt.border_width_pt > 0:
            shape.line.color.rgb = _hex_to_rgb(fmt.border_color_hex)
            shape.line.width = Pt(fmt.border_width_pt)

    except Exception as e:
        sid = getattr(shape, "shape_id", "?")
        print(f"    [경고] 서식 적용 실패 shape_id={sid}: {e}")
