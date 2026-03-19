"""
PPTX 통합 및 서식 통일 도구
사용법: python main.py [config.yaml 경로]
"""

import sys
from pathlib import Path
from pptx import Presentation

from config_loader import load_config, FormatConfig, SlideConfig
from slide_copier import copy_slide
from formatter import apply_format


def _set_dimensions(dest_prs: Presentation, src_prs: Presentation):
    dest_prs.slide_width = src_prs.slide_width
    dest_prs.slide_height = src_prs.slide_height


def _process_slide(
    src_prs: Presentation,
    slide_cfg: SlideConfig,
    fmt: FormatConfig,
    dest_prs: Presentation,
):
    slide_idx = slide_cfg.slide_index - 1  # 0-based 변환

    if slide_idx < 0 or slide_idx >= len(src_prs.slides):
        print(f"    [경고] 슬라이드 {slide_cfg.slide_index} 가 없습니다 — 건너뜀")
        return

    src_slide = src_prs.slides[slide_idx]
    dest_slide = copy_slide(src_slide, src_prs, dest_prs)

    include_ids = {s.shape_id for s in slide_cfg.shapes if s.include}
    exclude_ids = {s.shape_id for s in slide_cfg.shapes if not s.include}

    # include: false 도형 제거
    for shape in list(dest_slide.shapes):
        if shape.shape_id in exclude_ids:
            shape._element.getparent().remove(shape._element)

    # include: true 도형에 서식 적용
    for shape in dest_slide.shapes:
        if shape.shape_id in include_ids:
            apply_format(shape, fmt)

    print(
        f"    슬라이드 {slide_cfg.slide_index}: "
        f"포함 {len(include_ids)}개, 제외 {len(exclude_ids)}개"
    )


def run(config_path: str):
    print(f"설정 파일 로드: {config_path}")
    config = load_config(config_path)

    dest_prs = Presentation()
    dim_set = False

    for file_cfg in config.files:
        print(f"\n▶ {Path(file_cfg.path).name}")

        try:
            src_prs = Presentation(file_cfg.path)
        except Exception as e:
            print(f"  [오류] 파일 열기 실패: {e}")
            continue

        # 첫 번째 파일의 슬라이드 크기를 출력물에 적용
        if not dim_set:
            _set_dimensions(dest_prs, src_prs)
            dim_set = True

        for slide_cfg in file_cfg.slides:
            _process_slide(src_prs, slide_cfg, config.format, dest_prs)

    output_path = Path(config.output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    dest_prs.save(str(output_path))
    print(f"\n저장 완료: {output_path}")


if __name__ == "__main__":
    cfg_path = sys.argv[1] if len(sys.argv) > 1 else "config.yaml"
    run(cfg_path)
